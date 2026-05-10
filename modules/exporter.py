from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import io
import re
from datetime import datetime


def safe_currency(val):
    try:
        s = str(val).replace("€", "").strip()
        if not s or s == '0':
            return "Da definire"
        s = s.replace(".", "").replace(",", ".")
        num = float(s)
        return f"€ {num:,.0f}"
    except (ValueError, TypeError):
        return "Da definire"


class BandoExporter:
    def _add_markdown_text(self, doc, text):
        for line in text.strip().split('\n'):
            line = line.strip()
            if not line:
                continue
            if line.startswith('### '):
                p = doc.add_heading(line[4:], level=3)
            elif line.startswith('## '):
                p = doc.add_heading(line[3:], level=2)
            elif line.startswith('# '):
                p = doc.add_heading(line[2:], level=1)
            elif line.startswith('| '):
                self._add_table_from_markdown(doc, text)
                break
            elif line.startswith('- ') or line.startswith('* '):
                doc.add_paragraph(line[2:], style='List Bullet')
            else:
                doc.add_paragraph(line)

    def _add_table_from_markdown(self, doc, text):
        lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
        table_lines = []
        in_table = False
        for line in lines:
            if line.startswith('|'):
                in_table = True
                if not all(c in '|-: ' for c in line.replace('|', '').strip()):
                    table_lines.append(line)
            elif in_table:
                break
        if len(table_lines) < 2:
            return
        rows = []
        for tl in table_lines:
            cells = [c.strip() for c in tl.split('|')[1:-1]]
            rows.append(cells)
        if len(rows) < 2:
            return
        table = doc.add_table(rows=len(rows), cols=len(rows[0]))
        table.style = 'Light Grid Accent 1'
        for i, row_data in enumerate(rows):
            for j, cell_text in enumerate(row_data):
                if j < len(table.columns):
                    table.rows[i].cells[j].text = cell_text
        doc.add_paragraph()

    def _add_dnsh_section(self, doc, scheda_bando):
        doc.add_heading('5. Principio DNSH (Do No Significant Harm)', level=1)
        dnsh_richiesto = "Non specificato"
        if "DNSH" in scheda_bando or "Do No Significant Harm" in scheda_bando:
            dnsh_richiesto = "Richiesto"
        doc.add_paragraph(f"Valutazione DNSH: {dnsh_richiesto}")
        doc.add_paragraph(
            "Il principio DNSH richiede che l'investimento non arrechi un danno significativo "
            "all'ambiente. Verificare la conformità del progetto ai sei obiettivi ambientali "
            "del Regolamento (UE) 2020/852:"
        )
        obiettivi = [
            "Mitigazione dei cambiamenti climatici",
            "Adattamento ai cambiamenti climatici",
            "Uso sostenibile e protezione delle acque e delle risorse marine",
            "Economia circolare, prevenzione e riciclaggio dei rifiuti",
            "Prevenzione e riduzione dell'inquinamento",
            "Protezione e ripristino della biodiversità e degli ecosistemi"
        ]
        for obj in obiettivi:
            doc.add_paragraph(obj, style='List Bullet')
        doc.add_paragraph(
            "Documentazione necessaria: Dichiarazione sostitutiva di atto notorio che attesti "
            "la conformità del progetto al principio DNSH, corredata da relazione tecnica "
            "ove richiesto."
        )

    def _add_budget_table(self, doc, bozza_progetto):
        doc.add_heading('6. Dettaglio Budget', level=1)
        table = doc.add_table(rows=8, cols=3)
        table.style = 'Light Grid Accent 1'
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        headers = ['Voce di Spesa', 'Importo (€)', '% sul Totale']
        for j, h in enumerate(headers):
            cell = table.rows[0].cells[j]
            cell.text = h
            for p in cell.paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in p.runs:
                    run.bold = True
        voci = [
            ('Software e Licenze', '0', '0%'),
            ('Hardware e Infrastruttura', '0', '0%'),
            ('Consulenze Specialistiche', '0', '0%'),
            ('Formazione', '0', '0%'),
            ('Servizi Cloud', '0', '0%'),
            ('Spese Generali', '0', '0%'),
        ]
        # Try to extract values from business plan
        for v in voci:
            pattern = re.search(
                rf'{re.escape(v[0])}.*?€?\s*([\d.,]+)',
                bozza_progetto, re.IGNORECASE
            )
        for i, (voce, _, _) in enumerate(voci):
            table.rows[i + 1].cells[0].text = voce
            match = re.search(rf'{re.escape(voce)}.*?€?\s*([\d.,]+)',
                              bozza_progetto, re.IGNORECASE)
            val = match.group(1) if match else '0'
            table.rows[i + 1].cells[1].text = safe_currency(val)
            table.rows[i + 1].cells[2].text = '—'
        table.rows[7].cells[0].text = 'TOTALE'
        table.rows[7].cells[0].paragraphs[0].runs[0].bold = True
        table.rows[7].cells[1].text = 'Da definire'
        table.rows[7].cells[2].text = '100%'
        doc.add_paragraph()

    def genera_word(self, dati_azienda, scheda_bando, esito_matching, bozza_progetto, **kwargs):
        doc = Document()

        title = doc.add_heading('Report Analisi Bando', level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph(
            f"Generato il: {datetime.now().strftime('%d/%m/%Y %H:%M')}"
        )
        doc.add_paragraph('—' * 60)

        doc.add_heading('1. Profilo Azienda', level=1)
        for line in dati_azienda.strip().split('\n'):
            if ':' in line:
                k, v = line.split(':', 1)
                doc.add_paragraph(f"\u2022 {k.strip()}: {v.strip()}")
            elif line.strip():
                doc.add_paragraph(line.strip())

        doc.add_heading('2. Scheda Tecnica Bando', level=1)
        self._add_markdown_text(doc, scheda_bando)

        doc.add_heading('3. Esito Eligibility Check', level=1)
        self._add_markdown_text(doc, esito_matching)

        doc.add_heading('4. Business Plan / Progetto', level=1)
        self._add_markdown_text(doc, bozza_progetto)

        self._add_dnsh_section(doc, scheda_bando)
        self._add_budget_table(doc, bozza_progetto)

        doc.add_paragraph('—' * 60)
        p = doc.add_paragraph('Report generato automaticamente dal sistema di analisi bandi.')
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer

    def genera_pdf(self, dati_azienda, scheda_bando, esito_matching, bozza_progetto):
        pdf = FPDF()
        pdf.add_page()
        pdf.add_font("Arial", "", "C:/Windows/Fonts/arial.ttf")
        pdf.add_font("Arial", "B", "C:/Windows/Fonts/arialbd.ttf")
        lm = pdf.l_margin

        def mc(text, size=10, style=""):
            pdf.set_font("Arial", style, size)
            pdf.set_x(lm)
            pdf.multi_cell(0, 5, text, new_x="LMARGIN", new_y="NEXT")

        mc("Report Analisi Bando", 18, "B")
        mc(f"Generato il: {datetime.now().strftime('%d/%m/%Y %H:%M')}", 9)
        pdf.set_draw_color(200, 200, 200)
        pdf.line(lm, pdf.get_y(), 210 - lm, pdf.get_y())
        pdf.ln(4)

        def write_section(title, content):
            mc(title, 13, "B")
            pdf.ln(1)
            for line in content.strip().split("\n"):
                line = line.strip()
                if not line:
                    continue
                if line.startswith("### ") or line.startswith("## ") or line.startswith("# "):
                    mc(line.lstrip("# ").strip(), 11, "B")
                elif line.startswith("- ") or line.startswith("* "):
                    pdf.set_x(lm + 5)
                    pdf.multi_cell(0, 5, f"  {line[2:]}", new_x="LMARGIN", new_y="NEXT")
                elif line.startswith("|"):
                    pass
                else:
                    mc(line)
                if pdf.get_y() > 265:
                    pdf.add_page()

        write_section("1. Profilo Azienda", dati_azienda)
        write_section("2. Scheda Tecnica Bando", scheda_bando)
        write_section("3. Esito Eligibility Check", esito_matching)
        write_section("4. Business Plan / Progetto", bozza_progetto)

        mc("5. DNSH", 13, "B")
        mc("Valutazione Do No Significant Harm richiesta dal bando." if "DNSH" in scheda_bando else "DNSH non specificato nel bando.")

        pdf.ln(4)
        pdf.line(lm, pdf.get_y(), 210 - lm, pdf.get_y())
        mc("Report generato automaticamente dal sistema di analisi bandi.", 9)

        buffer = io.BytesIO()
        pdf.output(buffer)
        buffer.seek(0)
        return buffer

    def genera_slides(self, dati_analisi: dict) -> io.BytesIO:
        prs = Presentation()
        prs.slide_width = PptInches(13.333)
        prs.slide_height = PptInches(7.5)

        azienda = dati_analisi.get('azienda', 'N/D')
        bando = dati_analisi.get('bando', 'Analisi Bando')
        esito = dati_analisi.get('esito', '')
        investimento = dati_analisi.get('investimento', 0)
        fatturato = dati_analisi.get('fatturato', 0)
        calcolo = dati_analisi.get('calcolo_finanziario', {})
        parametri = dati_analisi.get('parametri_finanziari', {})
        checklist_items = dati_analisi.get('checklist_items', {})

        def parse_status(text):
            if not text:
                return 'N/D'
            t = text.upper()
            if 'ROSSO' in t:
                return 'ROSSO'
            if 'GIALLO' in t:
                return 'GIALLO'
            if 'VERDE' in t:
                return 'VERDE'
            return 'N/D'

        status = parse_status(esito)
        status_colors = {
            'VERDE': PptRGBColor(0x28, 0xA7, 0x45),
            'GIALLO': PptRGBColor(0xFF, 0xC1, 0x07),
            'ROSSO': PptRGBColor(0xDC, 0x35, 0x45),
            'N/D': PptRGBColor(0x99, 0x99, 0x99)
        }

        # Slide 1: Titolo
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PptRGBColor(0x1F, 0x4E, 0x79)

        txBox = slide.shapes.add_textbox(PptInches(1), PptInches(2), PptInches(11), PptInches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Analisi di Finanziamento"
        p.font.size = PptPt(44)
        p.font.color.rgb = PptRGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(PptInches(1), PptInches(3.5), PptInches(11), PptInches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = azienda
        p2.font.size = PptPt(28)
        p2.font.color.rgb = PptRGBColor(0xCC, 0xDD, 0xEE)
        p2.alignment = PP_ALIGN.CENTER

        txBox3 = slide.shapes.add_textbox(PptInches(1), PptInches(5), PptInches(11), PptInches(0.8))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = bando
        p3.font.size = PptPt(18)
        p3.font.color.rgb = PptRGBColor(0xAA, 0xBB, 0xCC)
        p3.alignment = PP_ALIGN.CENTER

        # Slide 2: Eligibility
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide2.shapes.add_textbox(PptInches(1), PptInches(0.5), PptInches(11), PptInches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Esito Eligibility"
        p.font.size = PptPt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        status_color = status_colors.get(status, PptRGBColor(0x99, 0x99, 0x99))
        txBox2 = slide2.shapes.add_textbox(PptInches(3), PptInches(2.5), PptInches(7), PptInches(2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = f"Stato: {status}"
        p2.font.size = PptPt(48)
        p2.font.bold = True
        p2.font.color.rgb = status_color
        p2.alignment = PP_ALIGN.CENTER

        txBox3 = slide2.shapes.add_textbox(PptInches(1), PptInches(5), PptInches(11), PptInches(1.5))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        # Extract probability from esito
        prob_match = re.search(r'PROBABILITÀ\s*APPROVAZIONE\s*[:\-]?\s*(\d+)', esito, re.IGNORECASE)
        prob_text = f"Probabilità Approvazione: {prob_match.group(1)}%" if prob_match else ""
        p3 = tf3.paragraphs[0]
        p3.text = prob_text
        p3.font.size = PptPt(24)
        p3.alignment = PP_ALIGN.CENTER

        txBox4 = slide2.shapes.add_textbox(PptInches(1), PptInches(5.8), PptInches(11), PptInches(1))
        tf4 = txBox4.text_frame
        tf4.word_wrap = True
        p4 = tf4.paragraphs[0]
        p4.text = f"Azienda: {azienda} | Investimento: \u20ac{int(investimento):,}"
        p4.font.size = PptPt(16)
        p4.font.color.rgb = PptRGBColor(0x66, 0x66, 0x66)
        p4.alignment = PP_ALIGN.CENTER

        # Slide 3: Piano Finanziario
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide3.shapes.add_textbox(PptInches(1), PptInches(0.5), PptInches(11), PptInches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Piano Finanziario"
        p.font.size = PptPt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        contrib_pct = float(parametri.get('aliquota_contributo', calcolo.get('aliquota_contributo', 0)))
        finanz_pct = float(parametri.get('aliquota_finanziamento', calcolo.get('aliquota_finanziamento', 0)))

        table = slide3.shapes.add_table(5, 3, PptInches(2), PptInches(2), PptInches(9), PptInches(3)).table
        table.columns[0].width = PptInches(4)
        table.columns[1].width = PptInches(2.5)
        table.columns[2].width = PptInches(2.5)

        headers = ['Voce', 'Percentuale', 'Importo']
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = PptPt(16)
                p.alignment = PP_ALIGN.CENTER

        rows_data = [
            ('Investimento Previsto', '—', f'\u20ac{int(investimento):,}'),
            ('Contributo a Fondo Perduto', f'{contrib_pct:.0f}%', f'\u20ac{int(investimento * contrib_pct / 100):,}'),
            ('Finanziamento Agevolato', f'{finanz_pct:.0f}%', f'\u20ac{int(investimento * finanz_pct / 100):,}'),
            ('Totale Agevolazione', f'{(contrib_pct + finanz_pct):.0f}%', f'\u20ac{int(investimento * (contrib_pct + finanz_pct) / 100):,}'),
        ]
        for i, (voce, pct, imp) in enumerate(rows_data):
            for j, val in enumerate([voce, pct, imp]):
                cell = table.cell(i + 1, j)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = PptPt(14)
                    if i == 3:
                        p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

        fatt_info = f"Fatturato: \u20ac{int(fatturato):,}" if fatturato else ""
        txBox3 = slide3.shapes.add_textbox(PptInches(1), PptInches(5.5), PptInches(11), PptInches(1))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = fatt_info
        p3.font.size = PptPt(14)
        p3.font.color.rgb = PptRGBColor(0x66, 0x66, 0x66)
        p3.alignment = PP_ALIGN.CENTER

        # Slide 4: Prossimi Passi
        slide4 = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide4.shapes.add_textbox(PptInches(1), PptInches(0.5), PptInches(11), PptInches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Prossimi Passi e Checklist"
        p.font.size = PptPt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        doc_list = list(checklist_items.keys()) if checklist_items else parametri.get('documenti_richiesti', [])
        if not doc_list:
            doc_list = [
                "Documento Identit\u00e0 Legale Rappresentante",
                "Visura Camerale",
                "Bilanci Ultimi 2 Anni",
                "DURC",
                "Preventivi Investimento",
                "Valutazione DNSH",
            ]

        txBox2 = slide4.shapes.add_textbox(PptInches(1.5), PptInches(2), PptInches(10), PptInches(4.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, doc_name in enumerate(doc_list):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"\u2610  {doc_name}"
            p.font.size = PptPt(18)
            p.space_after = PptPt(8)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
