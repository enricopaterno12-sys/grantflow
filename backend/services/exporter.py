import io
from typing import Optional
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from fpdf import FPDF
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment


# ── HELPERS ──

EMERALD = (16, 185, 129)
DARK_BG = (10, 10, 10)
WHITE = (255, 255, 255)
GRAY = (156, 163, 175)
DARK_GRAY = (30, 30, 30)


# ── PDF (fpdf2) ──

class GrantFlowPDF(FPDF):
    def header(self):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(*EMERALD)
        self.cell(0, 8, "GrantFlow AI", align="L")
        self.ln(12)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(*GRAY)
        self.cell(0, 10, f"Pagina {self.page_no()}/{{nb}}", align="C")


def genera_pdf(
    nome_bando: str,
    ente: str,
    azienda: dict,
    eligibility: str,
    analisi: str,
    business_plan: str,
    parametri: dict,
    calcolo: dict,
) -> bytes:
    pdf = GrantFlowPDF()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)

    # Cover
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 28)
    pdf.set_text_color(*EMERALD)
    pdf.cell(0, 20, "GrantFlow AI", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 14)
    pdf.set_text_color(*WHITE)
    pdf.cell(0, 10, "Report di Analisi Agevolazioni", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(10)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*GRAY)
    pdf.cell(0, 8, f"Bando: {nome_bando}", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 8, f"Ente: {ente}", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 8, f"Azienda: {azienda.get('ragione_sociale', 'N/D')}", align="C", new_x="LMARGIN", new_y="NEXT")

    # Eligibility
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_text_color(*EMERALD)
    pdf.cell(0, 12, "Report Eligibility", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(*WHITE)
    pdf.multi_cell(0, 6, eligibility)

    # Analisi Tecnica
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_text_color(*EMERALD)
    pdf.cell(0, 12, "Analisi Tecnica", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(*WHITE)
    pdf.multi_cell(0, 6, analisi)

    # Key Financials
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_text_color(*EMERALD)
    pdf.cell(0, 12, "Dati Finanziari", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(*WHITE)
    items = [
        f"Investimento: EUR{calcolo.get('investimento_effettivo', 0):,.0f}",
        f"Contributo: EUR{calcolo.get('contributo', 0):,.0f}",
        f"Finanziamento: EUR{calcolo.get('finanziamento', 0):,.0f}",
        f"DSCR: {calcolo.get('dscr', 'N/D')}",
        f"VAN: EUR{calcolo.get('van', 0):,.0f}",
        f"IRR: {calcolo.get('irr', 0)}%",
    ]
    for item in items:
        pdf.cell(0, 8, item, new_x="LMARGIN", new_y="NEXT")

    # Business Plan
    if business_plan:
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.set_text_color(*EMERALD)
        pdf.cell(0, 12, "Business Plan", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*WHITE)
        pdf.multi_cell(0, 6, business_plan[:3000])

    return pdf.output(dest="S").encode("latin-1", errors="replace")


# ── DOCX ──

def genera_docx(
    nome_bando: str,
    ente: str,
    azienda: dict,
    eligibility: str,
    analisi: str,
    business_plan: str,
    parametri: dict,
    calcolo: dict,
    checklist: list,
) -> bytes:
    doc = Document()

    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(*DARK_GRAY)

    # Title
    title = doc.add_heading("GrantFlow AI", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.color.rgb = RGBColor(*EMERALD)

    doc.add_heading(f"Dossier Tecnico — {azienda.get('ragione_sociale', 'N/D')}", level=1)
    doc.add_paragraph(f"Bando: {nome_bando}  |  Ente: {ente}")

    # 1. Sintesi
    doc.add_heading("1. Sintesi del Progetto", level=2)
    doc.add_paragraph(
        f"Il presente dossier illustra il progetto agevolato di {azienda.get('ragione_sociale', 'N/D')} "
        f"con investimento di EUR{calcolo.get('investimento_effettivo', 0):,.0f}, "
        f"contributo EUR{calcolo.get('contributo', 0):,.0f} e "
        f"finanziamento EUR{calcolo.get('finanziamento', 0):,.0f}."
    )

    # 2. Coerenza
    doc.add_heading("2. Coerenza con il Bando", level=2)
    doc.add_paragraph(f"ATECO azienda: {azienda.get('ateco', 'N/D')}")
    doc.add_paragraph(f"Regione: {azienda.get('regione', 'N/D')}")
    doc.add_paragraph(f"Dimensione: {azienda.get('dimensione', 'N/D')}")

    # 3. Eligibility
    doc.add_heading("3. Verifica Eligibility", level=2)
    for line in eligibility.split("\n"):
        if line.strip():
            doc.add_paragraph(line.strip())

    # 4. Piano Investimenti
    doc.add_heading("4. Piano Investimenti", level=2)
    table = doc.add_table(rows=4, cols=2, style="Light Grid Accent 1")
    data = [
        ("Investimento", f"EUR{calcolo.get('investimento_effettivo', 0):,.0f}"),
        ("Contributo", f"EUR{calcolo.get('contributo', 0):,.0f}"),
        ("Finanziamento", f"EUR{calcolo.get('finanziamento', 0):,.0f}"),
        ("Totale Agevolabile", f"EUR{calcolo.get('totale_agevolabile', 0):,.0f}"),
    ]
    for i, (k, v) in enumerate(data):
        table.rows[i].cells[0].text = k
        table.rows[i].cells[1].text = v

    # 5. Cronoprogramma
    doc.add_heading("5. Cronoprogramma", level=2)
    crono = doc.add_table(rows=4, cols=3, style="Light Grid Accent 1")
    for i, (fase, durata, periodo) in enumerate([
        ("Avvio progetto", "2 mesi", "Mese 1-2"),
        ("Sviluppo/Acquisti", "6 mesi", "Mese 3-8"),
        ("Collaudo", "2 mesi", "Mese 9-10"),
        ("Rendicontazione", "2 mesi", "Mese 11-12"),
    ]):
        crono.rows[i].cells[0].text = fase
        crono.rows[i].cells[1].text = durata
        crono.rows[i].cells[2].text = periodo

    # 6. Proiezioni
    doc.add_heading("6. Proiezioni Finanziarie", level=2)
    bp = calcolo
    doc.add_paragraph(f"DSCR: {bp.get('dscr', 'N/D')}  |  VAN: EUR{bp.get('van', 0):,.0f}  |  IRR: {bp.get('irr', 0)}%  |  Payback: {bp.get('payback_anni', 0)} anni")

    if business_plan:
        doc.add_heading("7. Business Plan", level=2)
        doc.add_paragraph(business_plan[:2000])

    # 8. Checklist
    doc.add_heading("8. Documenti Necessari", level=2)
    for item in checklist:
        doc.add_paragraph(f"[ ] {item.get('nome', '')}", style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PPTX ──

def genera_pptx(
    azienda: dict,
    calcolo: dict,
    eligibility: str,
    nome_bando: str,
    analisi: str,
) -> bytes:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    def add_slide():
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PptRGB(*DARK_BG)
        return slide

    # 1. Title
    s = add_slide()
    s.shapes.title.text = "GrantFlow AI"
    s.shapes.add_textbox(PptInches(1), PptInches(2), PptInches(11), PptInches(1)).text_frame.text = \
        f"Pitch: {azienda.get('ragione_sociale', 'N/D')}"
    s.shapes.add_textbox(PptInches(1), PptInches(3.5), PptInches(11), PptInches(0.5)).text_frame.text = \
        f"Bando: {nome_bando}"

    # 2. Overview
    s = add_slide()
    txbox = s.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(12), PptInches(1))
    txbox.text_frame.text = "Overview Progetto"
    lines = [
        f"Investimento: EUR{calcolo.get('investimento_effettivo', 0):,.0f}",
        f"Contributo: EUR{calcolo.get('contributo', 0):,.0f}",
        f"Finanziamento: EUR{calcolo.get('finanziamento', 0):,.0f}",
        f"DSCR: {calcolo.get('dscr', 'N/D')}  |  VAN: EUR{calcolo.get('van', 0):,.0f}  |  IRR: {calcolo.get('irr', 0)}%",
    ]
    txbox = s.shapes.add_textbox(PptInches(0.5), PptInches(1.5), PptInches(12), PptInches(4))
    txbox.text_frame.text = "\n".join(lines)

    # 3. Eligibility
    s = add_slide()
    txbox = s.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(12), PptInches(1))
    txbox.text_frame.text = "Eligibility Checks"
    txbox = s.shapes.add_textbox(PptInches(0.5), PptInches(1.5), PptInches(12), PptInches(5))
    lines = [line.strip() for line in eligibility.split("\n") if line.strip()][:15]
    txbox.text_frame.text = "\n".join(lines)

    # 4. Financials
    s = add_slide()
    txbox = s.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(12), PptInches(1))
    txbox.text_frame.text = "Key Financials"
    txbox = s.shapes.add_textbox(PptInches(0.5), PptInches(1.5), PptInches(12), PptInches(4))
    txbox.text_frame.text = f"""DSCR: {calcolo.get('dscr', 'N/D')}
VAN: EUR{calcolo.get('van', 0):,.0f}
IRR: {calcolo.get('irr', 0)}%
Payback: {calcolo.get('payback_anni', 0)} anni
Contributo: EUR{calcolo.get('contributo', 0):,.0f}
Finanziamento: EUR{calcolo.get('finanziamento', 0):,.0f}"""

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── XLSX ──

def genera_xlsx(
    calcolo: dict,
    azienda: dict,
    parametri: dict,
    nome_bando: str,
) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Business Plan"

    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="10B981", end_color="10B981", fill_type="solid")

    # Section 1: Dati Generali
    ws.append(["GrantFlow AI - Business Plan"])
    ws.merge_cells("A1:D1")
    ws.append([])
    ws.append(["Azienda", azienda.get("ragione_sociale", "")])
    ws.append(["Bando", nome_bando])
    ws.append(["ATECO", azienda.get("ateco", "")])
    ws.append(["Investimento", calcolo.get("investimento_effettivo", 0)])
    ws.append([])

    # Section 2: Key Metrics
    ws.append(["INDICATORI"])
    for cell in ws[ws.max_row]:
        cell.font = header_font
        cell.fill = header_fill
    ws.append(["DSCR", calcolo.get("dscr", "")])
    ws.append(["VAN", calcolo.get("van", "")])
    ws.append(["IRR (%)", calcolo.get("irr", "")])
    ws.append(["Payback (anni)", calcolo.get("payback_anni", "")])
    ws.append(["Contributo", calcolo.get("contributo", "")])
    ws.append(["Finanziamento", calcolo.get("finanziamento", "")])
    ws.append([])

    # Section 3: Cashflow
    ws.append(["PROIEZIONI CASHLOW"])
    for cell in ws[ws.max_row]:
        cell.font = header_font
        cell.fill = header_fill
    ws.append(["Anno", "Ricavi", "Costi", "Netto"])
    for row in calcolo.get("cashflow", []):
        ws.append([row["anno"], row["ricavi"], row["costi"], row["netto"]])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
